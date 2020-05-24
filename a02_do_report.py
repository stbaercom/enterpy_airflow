import os
import sqlite3
import ftplib

import pandas as pd
from pptx import Presentation
from pptx.util import Inches


def read_db(fn):
    con = sqlite3.connect(fn)
    df = pd.read_sql("select * from report;", con)
    return df

def write_pptx(df,fn):
    prs = Presentation()

    df = df.sort_values(by=["total"],ascending=False)[:5]

    producs = df['product'].to_list()
    totals = df['total'].to_list()

    # Cover
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Report"
    subtitle.text = "Produkte Verkaufszahlen"

    # Report
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Top 5 Produkte"

    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(1.5)
    shape = slide.shapes.add_table(rows=5, cols=2, left=x, top=y, width=cx, height=cy)
    table = shape.table

    for i,(prod,tot) in enumerate(zip(producs,totals)):
        table.cell(i,0).text = prod
        table.cell(i,1).text = f"{tot:.2f}"
    prs.save(fn)

def upload_pptx(fp):
    with ftplib.FTP() as session:
        session.connect('localhost',2121)
        session.login('ftp', 'ftp')
        with open(fp,mode="rb") as inf:
            fn = os.path.basename(fp)
            session.storbinary(f'STOR {fn}', inf)

def main():
    report_path = "data/report.pptx"
    df = read_db("data/report.db")
    write_pptx(df, report_path)
    upload_pptx(report_path)


if __name__ == '__main__':
    main()

